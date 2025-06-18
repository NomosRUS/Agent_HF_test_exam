import os
import gradio as gr
import requests
import inspect
import pandas as pd
from transformers import pipeline
from urllib.parse import quote_plus
import re
from typing import List, Dict, Any
from pptx import Presentation
from PyPDF2 import PdfReader

#to delete futher
current_question = ""
current_answer = ""
def get_status():
    return current_question, current_answer


# (Keep Constants as is)
# --- Constants ---
DEFAULT_API_URL = "https://agents-course-unit4-scoring.hf.space"

# --- Basic Agent Definition ---
# ----- THIS IS WERE YOU CAN BUILD WHAT YOU WANT ------
class BasicAgent:
    SYSTEM_PROMPT = (
        "You are a general AI assistant. I will ask you a question."
        "Report your thoughts, and finish your final answer with the following template: "
        "FINAL ANSWER: {Answer}" 
        "YOUR FINAL ANSWER should be a number OR as few words as possible OR a comma separated list of numbers and/or strings. "
        "If you are asked for a number, don't use comma to write your number neither use units such as $ or percent sign unless specified otherwise. "
        "If you are asked for a string, don't use articles, neither abbreviations, and write the digits in plain text unless specified otherwise. "
        "If you are asked for a comma separated list, apply the above rules depending of whether the element to be put in the list is a number or a string."
        "Always end your output exactly with FINAL ANSWER: <Answer> and do not add any text after that."
    )
    
    #def __init__(self, model_name="google/flan-t5-base"):
    def __init__(self, model_name="Qwen/Qwen3-14B"):
        # mistralai/Mistral-7B-Instruct-v0.3
    
    
        token = os.getenv("HF_API_TOKEN")
        self.generator = pipeline("text-generation", model=model_name, token=token)
        self.memory: List[str] = []

    def _plan(self, question: str) -> List[Dict[str, Any]]:
        """Create a simple plan consisting of tool steps."""
        steps: List[Dict[str, Any]] = []

        file_match = re.search(r"(\S+\.(?:pdf|xlsx|csv|pptx|txt))", question, re.IGNORECASE)
        if file_match:
            steps.append({"tool": "file", "path": file_match.group(1)})
            return steps

        expr_match = re.search(r"[\d\s\+\-\*/\.\(\)]+", question)
        has_op = any(op in question for op in ["+", "-", "*", "/"])

        if expr_match and has_op:
            expression = expr_match.group(0)
            # if the question also implies web lookup, gather info first
            if re.search(r"\b(search|lookup|population|when|who|what)\b", question.lower()):
                steps.append({"tool": "web", "query": question})
            steps.append({"tool": "calculator", "expression": expression})
            return steps

        # default to web search
        steps.append({"tool": "web", "query": question})
        return steps

    def _web_search(self, query: str) -> str:
        url = f"https://r.jina.ai/https://duckduckgo.com/html/?q={quote_plus(query)}"
        try:
            resp = requests.get(url, timeout=10)
            text = resp.text
            for line in text.splitlines():
                if line.startswith("[") and "](" in line:
                    return line
            return "No result found"
        except Exception as e:
            return f"web search error: {e}"

    def _execute_calculator(self, expression: str) -> str:
        try:
            result = eval(expression, {"__builtins__": {}}, {})
            return str(result)
        except Exception as e:
            return f"calc error: {e}"

    def _load_file(self, path: str) -> str:
        try:
            ext = os.path.splitext(path)[1].lower()
            if ext == ".pdf":
                reader = PdfReader(path)
                return "\n".join(page.extract_text() for page in reader.pages[:3])
            if ext in {".xlsx", ".xls"}:
                df = pd.read_excel(path)
                return df.to_csv(index=False)
            if ext == ".csv":
                df = pd.read_csv(path)
                return df.to_csv(index=False)
            if ext == ".pptx":
                prs = Presentation(path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts)
            if ext == ".txt":
                with open(path, "r", encoding="utf-8") as f:
                    return f.read()
        except Exception as e:
            return f"file load error: {e}"
        return "unsupported file"

    def __call__(self, question: str) -> str:

        self.memory.clear()
        self.memory.append(f"Question: {question}")
        plan = self._plan(question)
        self.memory.append(f"Plan: {plan}")

        for step in plan:
            action = step.get("tool")
            self.memory.append(f"Act: {action} -> {step}")
            if action == "calculator":
                observation = self._execute_calculator(step["expression"])
            elif action == "file":
                observation = self._load_file(step["path"])
            else:
                observation = self._web_search(step["query"])
            self.memory.append(f"Observation: {observation}")

        context = "\n".join(self.memory)
        prompt = f"{self.SYSTEM_PROMPT}\n{context}\nQuestion: {question}\nAnswer:"
        
        try:
            outputs = self.generator(prompt, max_new_tokens=128)
        except Exception as e:
            raise RuntimeError(f"generation failed: {e}") from e

        if outputs and isinstance(outputs, list):
            generated_text = outputs[0].get("generated_text", "")
        else:
            generated_text = str(outputs)

        if generated_text.startswith(prompt):
            generated_text = generated_text[len(prompt):].lstrip()

        if "FINAL ANSWER:" in generated_text:
            return generated_text.split("FINAL ANSWER:", 1)[1].strip()
        return generated_text.strip()

def run_and_submit_all( profile: gr.OAuthProfile | None):
    """
    Fetches all questions, runs the BasicAgent on them, submits all answers,
    and displays the results.
    """
    # --- Determine HF Space Runtime URL and Repo URL ---
    space_id = os.getenv("SPACE_ID") # Get the SPACE_ID for sending link to the code

    if profile:
        username= f"{profile.username}"
        print(f"User logged in: {username}")
    else:
        print("User not logged in.")
        return "Please Login to Hugging Face with the button.", None

    api_url = DEFAULT_API_URL
    questions_url = f"{api_url}/questions"
    submit_url = f"{api_url}/submit"
    
    #delete
    global current_question, current_answer
    current_question = ""
    current_answer = ""
    #delete
    
    # 1. Instantiate Agent ( modify this part to create your agent)
    try:
        agent = BasicAgent()
    except Exception as e:
        print(f"Error instantiating agent: {e}")
        return f"Error initializing agent: {e}", None
    # In the case of an app running as a hugging Face space, this link points toward your codebase ( usefull for others so please keep it public)
    agent_code = f"https://huggingface.co/spaces/{space_id}/tree/main"
    print(agent_code)

    # 2. Fetch Questions
    print(f"Fetching questions from: {questions_url}")
    try:
        response = requests.get(questions_url, timeout=15)
        response.raise_for_status()
        questions_data = response.json()
        if not questions_data:
             print("Fetched questions list is empty.")
             return "Fetched questions list is empty or invalid format.", None
        print(f"Fetched {len(questions_data)} questions.")
    except requests.exceptions.RequestException as e:
        print(f"Error fetching questions: {e}")
        return f"Error fetching questions: {e}", None
    except requests.exceptions.JSONDecodeError as e:
         print(f"Error decoding JSON response from questions endpoint: {e}")
         print(f"Response text: {response.text[:500]}")
         return f"Error decoding server response for questions: {e}", None
    except Exception as e:
        print(f"An unexpected error occurred fetching questions: {e}")
        return f"An unexpected error occurred fetching questions: {e}", None

    # 3. Run your Agent
    results_log = []
    answers_payload = []
    print(f"Running agent on {len(questions_data)} questions...")
    for item in questions_data:
        task_id = item.get("task_id")
        question_text = item.get("question")
        #del
        current_question = question_text
        current_answer = ""
        #del
        if not task_id or question_text is None:
            print(f"Skipping item with missing task_id or question: {item}")
            continue
        try:
            submitted_answer = agent(question_text)
            answers_payload.append({"task_id": task_id, "submitted_answer": submitted_answer})
            #del
            current_answer = submitted_answer
            #del
            results_log.append({"Task ID": task_id, "Question": question_text, "Submitted Answer": submitted_answer})
        except Exception as e:
             print(f"Error running agent on task {task_id}: {e}")
             results_log.append({"Task ID": task_id, "Question": question_text, "Submitted Answer": f"AGENT ERROR: {e}"})
             #del
             current_answer = f"AGENT ERROR: {e}"
             #del
    if not answers_payload:
        print("Agent did not produce any answers to submit.")
        return "Agent did not produce any answers to submit.", pd.DataFrame(results_log)

    # 4. Prepare Submission 
    submission_data = {"username": username.strip(), "agent_code": agent_code, "answers": answers_payload}
    status_update = f"Agent finished. Submitting {len(answers_payload)} answers for user '{username}'..."
    print(status_update)

    # 5. Submit
    print(f"Submitting {len(answers_payload)} answers to: {submit_url}")
    try:
        response = requests.post(submit_url, json=submission_data, timeout=60)
        response.raise_for_status()
        result_data = response.json()
        final_status = (
            f"Submission Successful!\n"
            f"User: {result_data.get('username')}\n"
            f"Overall Score: {result_data.get('score', 'N/A')}% "
            f"({result_data.get('correct_count', '?')}/{result_data.get('total_attempted', '?')} correct)\n"
            f"Message: {result_data.get('message', 'No message received.')}"
        )
        print("Submission successful.")
        results_df = pd.DataFrame(results_log)
        return final_status, results_df
    except requests.exceptions.HTTPError as e:
        error_detail = f"Server responded with status {e.response.status_code}."
        try:
            error_json = e.response.json()
            error_detail += f" Detail: {error_json.get('detail', e.response.text)}"
        except requests.exceptions.JSONDecodeError:
            error_detail += f" Response: {e.response.text[:500]}"
        status_message = f"Submission Failed: {error_detail}"
        print(status_message)
        results_df = pd.DataFrame(results_log)
        return status_message, results_df
    except requests.exceptions.Timeout:
        status_message = "Submission Failed: The request timed out."
        print(status_message)
        results_df = pd.DataFrame(results_log)
        return status_message, results_df
    except requests.exceptions.RequestException as e:
        status_message = f"Submission Failed: Network error - {e}"
        print(status_message)
        results_df = pd.DataFrame(results_log)
        return status_message, results_df
    except Exception as e:
        status_message = f"An unexpected error occurred during submission: {e}"
        print(status_message)
        results_df = pd.DataFrame(results_log)
        return status_message, results_df


# --- Build Gradio Interface using Blocks ---
with gr.Blocks() as demo:
    gr.Markdown("# Basic Agent Evaluation Runner")
    gr.Markdown(
        """
        **Instructions:**

        1.  Please clone this space, then modify the code to define your agent's logic, the tools, the necessary packages, etc ...
        2.  Log in to your Hugging Face account using the button below. This uses your HF username for submission.
        3.  Click 'Run Evaluation & Submit All Answers' to fetch questions, run your agent, submit answers, and see the score.

        ---
        **Disclaimers:**
        Once clicking on the "submit button, it can take quite some time ( this is the time for the agent to go through all the questions).
        This space provides a basic setup and is intentionally sub-optimal to encourage you to develop your own, more robust solution. For instance for the delay process of the submit button, a solution could be to cache the answers and submit in a seperate action or even to answer the questions in async.
        """
    )

    gr.LoginButton()

    run_button = gr.Button("Run Evaluation & Submit All Answers")

    status_output = gr.Textbox(label="Run Status / Submission Result", lines=5, interactive=False)
    # Removed max_rows=10 from DataFrame constructor
    results_table = gr.DataFrame(label="Questions and Agent Answers", wrap=True)
    #del
    current_question_box = gr.Textbox(label="Current Question")
    current_answer_box = gr.Textbox(label="Current Answer")
    #del
    run_button.click(
        fn=run_and_submit_all,
        outputs=[status_output, results_table]
    )
    #del
    status_timer = gr.Timer(1.0)
    status_timer.tick(fn=get_status, outputs=[current_question_box, current_answer_box])
    #del
if __name__ == "__main__":
    print("\n" + "-"*30 + " App Starting " + "-"*30)
    # Check for SPACE_HOST and SPACE_ID at startup for information
    space_host_startup = os.getenv("SPACE_HOST")
    space_id_startup = os.getenv("SPACE_ID") # Get SPACE_ID at startup

    if space_host_startup:
        print(f"✅ SPACE_HOST found: {space_host_startup}")
        print(f"   Runtime URL should be: https://{space_host_startup}.hf.space")
    else:
        print("ℹ️  SPACE_HOST environment variable not found (running locally?).")

    if space_id_startup: # Print repo URLs if SPACE_ID is found
        print(f"✅ SPACE_ID found: {space_id_startup}")
        print(f"   Repo URL: https://huggingface.co/spaces/{space_id_startup}")
        print(f"   Repo Tree URL: https://huggingface.co/spaces/{space_id_startup}/tree/main")
    else:
        print("ℹ️  SPACE_ID environment variable not found (running locally?). Repo URL cannot be determined.")

    print("-"*(60 + len(" App Starting ")) + "\n")

    print("Launching Gradio Interface for Basic Agent Evaluation...")
    demo.launch(debug=True, share=False)